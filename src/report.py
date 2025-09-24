import pandas as pd
from pptx import Presentation
from pptx.util import Inches

def export_report(reconciled: pd.DataFrame, forecast: pd.DataFrame, output: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Budget Reconciliation Report"

    # Merge for summary
    df = reconciled.merge(forecast, on="Category")

    rows, cols = df.shape
    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
    table.columns[0].width = Inches(2)

    # headers
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = col

    # data
    for i in range(rows):
        for j, val in enumerate(df.iloc[i]):
            table.cell(i+1, j).text = str(round(val, 2)) if isinstance(val, (int, float)) else str(val)

    prs.save(output)